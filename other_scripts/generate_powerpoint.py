#!/usr/bin/env python3
"""
Generate a PowerPoint presentation from sample questions with video frames.

Each question slide has:
  1. Question type as heading
  2. Question text with video name and inserted frames
  3. Answers in two columns, correct answer highlighted yellow
"""

import argparse
import json
import os
import random
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SECONDARY_TYPES = {
    "role_count_aggressor", "role_count_victim", "role_count_bystander",
    "compound_aggressor_victim_count", "compound_victim_bystander_count",
    "compound_action_location",
}

# Slide dimensions (standard widescreen 13.333 x 7.5 inches)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
YELLOW_HIGHLIGHT = RGBColor(0xFF, 0xFF, 0x00)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY_BG = RGBColor(0xF2, 0xF2, 0xF2)


def _format_type_name(question_type: str) -> str:
    """Convert question_type slug to a readable heading."""
    return question_type.replace("_", " ").title()


def _add_text_box(slide, left, top, width, height, text, font_size=12,
                  bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                  font_name="Calibri"):
    """Add a simple text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_answer_box(slide, left, top, width, height, number, text,
                    is_correct=False):
    """Add a single answer box, highlighted yellow if correct."""
    from pptx.oxml.ns import qn

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)

    if is_correct:
        solid = txBox.fill
        solid.solid()
        solid.fore_color.rgb = YELLOW_HIGHLIGHT
        txBox.line.color.rgb = RGBColor(0xCC, 0xCC, 0x00)
        txBox.line.width = Pt(1)
    else:
        solid = txBox.fill
        solid.solid()
        solid.fore_color.rgb = WHITE
        txBox.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        txBox.line.width = Pt(0.5)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{number}. {text}"
    run.font.size = Pt(10)
    run.font.name = "Calibri"
    run.font.color.rgb = DARK_GRAY
    if is_correct:
        run.font.bold = True
    p.alignment = PP_ALIGN.LEFT


def build_question_slide(prs, question, frames_dir, frames_per_video=2):
    """
    Build a single question slide.

    Layout:
      - Top: question type heading
      - Middle-left: video frames side by side with video name caption
      - Middle-right: question text
      - Bottom: answers in two columns, correct answer highlighted yellow
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    video_name = question["video_name"]
    q_type = question["question_type"]
    prompt = question["prompt"]
    answers = question["answers"]
    correct_idx = question["correct_index"]

    # --- 1. Heading: question type ---
    _add_text_box(
        slide,
        left=Inches(0.5), top=Inches(0.2),
        width=Inches(12), height=Inches(0.5),
        text=_format_type_name(q_type),
        font_size=24, bold=True, color=DARK_GRAY,
    )

    # --- 2. Video name ---
    _add_text_box(
        slide,
        left=Inches(0.5), top=Inches(0.7),
        width=Inches(12), height=Inches(0.3),
        text=f"Video: {video_name}",
        font_size=11, bold=False, color=RGBColor(0x66, 0x66, 0x66),
    )

    # --- 3. Video frames ---
    base_name = os.path.splitext(video_name)[0]
    frame_files = sorted(
        Path(frames_dir).glob(f"{base_name}_frame*.jpg")
    )[:frames_per_video]

    frame_left = Inches(0.5)
    frame_top = Inches(1.1)
    frame_width = Inches(2.5)
    frame_height = Inches(1.8)

    for i, frame_path in enumerate(frame_files):
        slide.shapes.add_picture(
            str(frame_path),
            frame_left + i * (frame_width + Inches(0.2)),
            frame_top,
            frame_width,
            frame_height,
        )

    # --- 4. Question text (to the right of frames, or below if no frames) ---
    if frame_files:
        q_left = frame_left + len(frame_files) * (frame_width + Inches(0.2)) + Inches(0.3)
        q_top = Inches(1.1)
        q_width = SLIDE_WIDTH - q_left - Inches(0.5)
    else:
        q_left = Inches(0.5)
        q_top = Inches(1.1)
        q_width = Inches(12)

    _add_text_box(
        slide,
        left=q_left, top=q_top,
        width=q_width, height=Inches(1.8),
        text=prompt,
        font_size=14, bold=False, color=BLACK,
    )

    # --- 5. Answers in two columns ---
    answers_top = Inches(3.2)
    col_width = Inches(5.8)
    col_gap = Inches(0.5)
    left_col_x = Inches(0.5)
    right_col_x = left_col_x + col_width + col_gap
    row_height = Inches(0.48)
    row_gap = Inches(0.08)

    mid = (len(answers) + 1) // 2  # left column gets the extra one if odd

    for i, answer in enumerate(answers):
        is_correct = (i == correct_idx)
        if i < mid:
            col_x = left_col_x
            row_idx = i
        else:
            col_x = right_col_x
            row_idx = i - mid

        box_top = answers_top + row_idx * (row_height + row_gap)
        _add_answer_box(
            slide, col_x, box_top, col_width, row_height,
            number=i + 1, text=answer, is_correct=is_correct,
        )


def main():
    parser = argparse.ArgumentParser(
        description="Generate a PowerPoint from sample questions with video frames"
    )
    parser.add_argument("questions", help="Path to generated questions JSON file")
    parser.add_argument("-o", "--output", default="sample_questions.pptx",
                        help="Output PPTX file (default: sample_questions.pptx)")
    parser.add_argument("-f", "--frames-dir", default="sample_frames",
                        help="Directory containing extracted frames (default: sample_frames)")
    parser.add_argument("-n", "--count", type=int, default=2,
                        help="Number of sample questions per category (default: 2)")
    parser.add_argument("--frames-per-video", type=int, default=8,
                        help="Max frames to show per video (default: 8)")
    parser.add_argument("-s", "--seed", type=int, default=None,
                        help="Random seed for reproducibility")
    args = parser.parse_args()

    with open(args.questions, "r", encoding="utf-8") as f:
        data = json.load(f)

    questions = []
    for video, qs in data.get("questions_by_video", {}).items():
        questions.extend(qs)

    by_type = {}
    for q in questions:
        qt = q["question_type"]
        if qt not in SECONDARY_TYPES:
            by_type.setdefault(qt, []).append(q)

    if args.seed is not None:
        random.seed(args.seed)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = title_slide.shapes.title
    if title_shape:
        title_shape.text = "Video Aggression Understanding"
    if 1 in title_slide.placeholders:
        title_slide.placeholders[1].text = "Sample Questions by Category"

    for qt in sorted(by_type.keys()):
        pool = by_type[qt]
        n = min(args.count, len(pool))
        samples = random.sample(pool, n)

        for q in samples:
            build_question_slide(prs, q, args.frames_dir, args.frames_per_video)

    prs.save(args.output)
    print(f"Saved {args.output} with {len(prs.slides) - 1} question slides")


if __name__ == "__main__":
    main()
